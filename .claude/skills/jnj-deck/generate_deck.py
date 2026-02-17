#!/usr/bin/env python3
"""
JnJ Branded PowerPoint Deck Generator
Reusable script — accepts a JSON slide definition via stdin or file argument.

Usage:
    echo '<json>' | python3 generate_deck.py --output my-deck.pptx
    python3 generate_deck.py --input slides.json --output my-deck.pptx

JSON Schema:
{
    "title": "Deck Title",
    "subtitle": "Subtitle text",
    "output": "filename.pptx",
    "slides": [
        {
            "type": "title",
            "title": "Presentation Title",
            "subtitle": "Optional subtitle"
        },
        {
            "type": "bullets",
            "title": "Slide Title",
            "bullets": ["Point 1", "Point 2"]
        },
        {
            "type": "metrics",
            "title": "Key Metrics",
            "metrics": [
                {"value": "$10B", "label": "Revenue"},
                {"value": "25%", "label": "Growth"}
            ]
        },
        {
            "type": "two_column",
            "title": "Comparison",
            "left_title": "Left Header",
            "left_bullets": ["A", "B"],
            "right_title": "Right Header",
            "right_bullets": ["C", "D"]
        },
        {
            "type": "section",
            "title": "Section Title"
        },
        {
            "type": "thankyou",
            "title": "Thank you",
            "subtitle": "Questions & Discussion"
        }
    ]
}
"""

import json
import sys
import os
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# === JnJ Brand Colors ===
JNJ_RED = RGBColor(213, 25, 0)
JNJ_DARK_RED = RGBColor(155, 6, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(26, 26, 26)
LIGHT_GRAY = RGBColor(242, 242, 242)
MEDIUM_GRAY = RGBColor(109, 110, 113)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Arial"


# ── Helper Functions ──

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_red_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=SLIDE_WIDTH, height=Inches(0.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = JNJ_RED
    bar.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(
        left=Inches(12.0), top=Inches(7.0),
        width=Inches(1.0), height=Inches(0.4),
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT


def add_text_box(slide, text, left, top, width, height, font_size, color,
                 bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        left=left, top=top, width=width, height=height,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_footnotes(slide, footnotes, bottom_y=Inches(6.6)):
    """Add numbered footnote hyperlinks at the bottom of a slide.

    Args:
        slide: The slide object.
        footnotes: List of dicts with "label" and "url" keys.
                   e.g. [{"label": "JnJ 2025 Annual Report", "url": "https://..."}]
        bottom_y: Vertical position for the footnote text box.
    """
    if not footnotes:
        return

    txBox = slide.shapes.add_textbox(
        left=Inches(0.75), top=bottom_y,
        width=Inches(11.8), height=Inches(0.8),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, fn in enumerate(footnotes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Add the number prefix as a plain run
        prefix_run = p.add_run()
        prefix_run.text = f"[{i + 1}] "
        prefix_run.font.size = Pt(9)
        prefix_run.font.color.rgb = MEDIUM_GRAY
        prefix_run.font.name = FONT_NAME

        # Add the label as a hyperlinked run
        link_run = p.add_run()
        link_run.text = fn.get("label", fn["url"])
        link_run.font.size = Pt(9)
        link_run.font.color.rgb = RGBColor(0, 102, 204)  # Blue hyperlink color
        link_run.font.name = FONT_NAME
        link_run.font.underline = True
        link_run.hyperlink.address = fn["url"]

        p.space_before = Pt(1)
        p.space_after = Pt(1)


def add_white_separator(slide, top):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(4.0), top=top,
        width=Inches(5.3), height=Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()


# ── Slide Builders ──

def build_title_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, JNJ_RED)
    add_text_box(slide, data["title"],
                 Inches(0.75), Inches(2.0), Inches(11.8), Inches(1.2),
                 Pt(44), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        add_text_box(slide, data["subtitle"],
                     Inches(0.75), Inches(3.4), Inches(11.8), Inches(1.0),
                     Pt(22), WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_white_separator(slide, Inches(4.5))
    return slide


def build_bullets_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_red_accent_bar(slide)
    add_text_box(slide, data["title"],
                 Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.2),
                 Pt(32), JNJ_RED, bold=True, alignment=PP_ALIGN.LEFT)

    has_footnotes = bool(data.get("footnotes"))
    body_height = Inches(4.2) if has_footnotes else Inches(5.0)

    txBox = slide.shapes.add_textbox(
        left=Inches(1.0), top=Inches(1.8),
        width=Inches(11.0), height=body_height,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(10)

    if has_footnotes:
        add_footnotes(slide, data["footnotes"])

    add_slide_number(slide, slide_num)
    return slide


def build_metrics_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_GRAY)
    add_red_accent_bar(slide)
    add_text_box(slide, data["title"],
                 Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.2),
                 Pt(32), JNJ_RED, bold=True, alignment=PP_ALIGN.LEFT)

    metrics = data["metrics"]
    box_count = len(metrics)
    total_width = 11.0
    box_width = (total_width - (box_count - 1) * 0.4) / box_count
    start_left = 1.0

    for i, m in enumerate(metrics):
        left = start_left + i * (box_width + 0.4)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(left), top=Inches(2.2),
            width=Inches(box_width), height=Inches(2.5),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = JNJ_RED
        box.line.fill.background()
        box.text_frame.word_wrap = True
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = box.text_frame.paragraphs[0]
        p.text = m["value"]
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME
        p.font.bold = True

        label_box = slide.shapes.add_textbox(
            left=Inches(left), top=Inches(4.9),
            width=Inches(box_width), height=Inches(0.6),
        )
        lp = label_box.text_frame.paragraphs[0]
        lp.text = m["label"]
        lp.font.size = Pt(14)
        lp.font.color.rgb = BLACK
        lp.font.name = FONT_NAME
        lp.alignment = PP_ALIGN.CENTER

    if data.get("footnotes"):
        add_footnotes(slide, data["footnotes"])

    add_slide_number(slide, slide_num)
    return slide


def build_two_column_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_red_accent_bar(slide)
    add_text_box(slide, data["title"],
                 Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.2),
                 Pt(32), JNJ_RED, bold=True, alignment=PP_ALIGN.LEFT)

    has_footnotes = bool(data.get("footnotes"))
    col_height = Inches(3.7) if has_footnotes else Inches(4.5)
    divider_height = Inches(4.0) if has_footnotes else Inches(4.8)

    # Left column header
    add_text_box(slide, data["left_title"],
                 Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.5),
                 Pt(20), JNJ_DARK_RED, bold=True)

    # Left column bullets
    lBox = slide.shapes.add_textbox(
        left=Inches(1.0), top=Inches(2.3),
        width=Inches(5.0), height=col_height,
    )
    ltf = lBox.text_frame
    ltf.word_wrap = True
    for i, b in enumerate(data["left_bullets"]):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(8)

    # Right column header
    add_text_box(slide, data["right_title"],
                 Inches(6.75), Inches(1.7), Inches(5.5), Inches(0.5),
                 Pt(20), JNJ_DARK_RED, bold=True)

    # Right column bullets
    rBox = slide.shapes.add_textbox(
        left=Inches(7.0), top=Inches(2.3),
        width=Inches(5.0), height=col_height,
    )
    rtf = rBox.text_frame
    rtf.word_wrap = True
    for i, b in enumerate(data["right_bullets"]):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"\u2022  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        p.font.name = FONT_NAME
        p.space_after = Pt(8)

    # Center divider
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(6.4), top=Inches(1.8),
        width=Inches(0.02), height=divider_height,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_GRAY
    line.line.fill.background()

    if has_footnotes:
        add_footnotes(slide, data["footnotes"])

    add_slide_number(slide, slide_num)
    return slide


def build_section_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, JNJ_RED)
    add_text_box(slide, data["title"],
                 Inches(0.75), Inches(2.8), Inches(11.8), Inches(1.2),
                 Pt(40), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_white_separator(slide, Inches(4.2))
    add_slide_number(slide, slide_num)
    return slide


def build_thankyou_slide(prs, data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, JNJ_RED)
    add_text_box(slide, data.get("title", "Thank you"),
                 Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.2),
                 Pt(44), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        add_text_box(slide, data["subtitle"],
                     Inches(0.75), Inches(3.8), Inches(11.8), Inches(1.0),
                     Pt(24), WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_white_separator(slide, Inches(5.0))
    return slide


# ── Slide type registry ──
BUILDERS = {
    "title": build_title_slide,
    "bullets": build_bullets_slide,
    "metrics": build_metrics_slide,
    "two_column": build_two_column_slide,
    "section": build_section_slide,
    "thankyou": build_thankyou_slide,
}


def main():
    parser = argparse.ArgumentParser(description="Generate a JnJ-branded PowerPoint deck from JSON")
    parser.add_argument("--input", "-i", help="Path to JSON file (default: read from stdin)")
    parser.add_argument("--output", "-o", help="Output .pptx filename (overrides JSON 'output' field)")
    args = parser.parse_args()

    # Read JSON input
    if args.input:
        with open(args.input, "r") as f:
            deck_data = json.load(f)
    else:
        deck_data = json.load(sys.stdin)

    # Determine output path
    output_path = args.output or deck_data.get("output", "deck.pptx")

    # Build presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_num = 0
    for slide_def in deck_data["slides"]:
        slide_type = slide_def.get("type", "bullets")
        builder = BUILDERS.get(slide_type)
        if not builder:
            print(f"Warning: Unknown slide type '{slide_type}', skipping.", file=sys.stderr)
            continue
        slide_num += 1
        builder(prs, slide_def, slide_num)

    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    prs.save(output_path)
    print(f"Deck saved to: {output_path}")
    print(f"Total slides: {slide_num}")


if __name__ == "__main__":
    main()
